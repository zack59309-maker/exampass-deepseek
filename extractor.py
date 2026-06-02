"""提取 PPTX/DOCX/PDF 文件的文本内容。"""

import os
import re

def extract_pptx(filepath):
    """Extract text from PPTX."""
    from pptx import Presentation
    prs = Presentation(filepath)
    text_parts = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            text_parts.append('\n'.join(slide_texts))
    return '\n\n---\n\n'.join(text_parts)


def extract_docx(filepath):
    """Extract text from DOCX."""
    from docx import Document
    doc = Document(filepath)
    text_parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            text_parts.append(t)
    # Also extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(' | '.join(cells))
        text_parts.append('\n'.join(rows))
    return '\n\n'.join(text_parts)


def extract_pdf(filepath):
    """Extract text from PDF."""
    import fitz
    doc = fitz.open(filepath)
    text_parts = []
    for page in doc:
        t = page.get_text().strip()
        if t:
            text_parts.append(t)
    return '\n\n'.join(text_parts)


def extract_file(filepath):
    """Auto-detect format and extract text."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.pptx':
        text = extract_pptx(filepath)
    elif ext == '.docx':
        text = extract_docx(filepath)
    elif ext == '.pdf':
        text = extract_pdf(filepath)
    else:
        raise ValueError(f"Unsupported format: {ext}")
    return text


def merge_texts(texts):
    """Merge multiple extracted texts with separators."""
    merged = []
    for i, t in enumerate(texts):
        merged.append(f"[文件 {i+1}]\n{t}")
    return '\n\n=====\n\n'.join(merged)
