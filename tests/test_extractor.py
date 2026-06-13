"""Smoke tests for extractor functions."""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import tempfile


def test_extract_file_raises_on_unsupported():
    """extract_file should raise ValueError for unsupported formats."""
    from extractor import extract_file
    try:
        extract_file("test.txt")
        assert False, "Should have raised ValueError"
    except ValueError as e:
        assert "Unsupported" in str(e)


def test_extract_file_pptx(tmp_path):
    """Create a minimal PPTX and extract text."""
    from extractor import extract_pptx
    from pptx import Presentation
    from pptx.util import Inches

    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Hello World Test"

    prs.save(str(pptx_path))

    text = extract_pptx(str(pptx_path))
    assert "Hello World Test" in text


def test_merge_texts():
    """merge_texts should combine texts with markers."""
    from extractor import merge_texts
    result = merge_texts(["alpha", "beta"])
    assert "[文件 1]" in result
    assert "[文件 2]" in result
    assert "alpha" in result
    assert "beta" in result
